#!/usr/bin/env python3
"""EndoCode: build .pptx from markdown. Slides: ## title, then lines or - bullets. Requires: pip install python-pptx"""
from __future__ import annotations

import json
import sys
from pathlib import Path


def parse_slides(md: str, fallback: str) -> list[dict]:
    slides: list[dict] = []
    cur: dict | None = None
    for line in md.splitlines():
        raw = line.rstrip()
        if raw.startswith("## "):
            if cur:
                slides.append(cur)
            cur = {"title": raw[3:].strip(), "bullets": []}
        elif raw.startswith("# ") and cur is None and not slides:
            cur = {"title": raw[2:].strip(), "bullets": []}
        elif cur is not None:
            t = raw.strip()
            if not t:
                continue
            if t.startswith(("- ", "* ")):
                cur["bullets"].append(t[2:].strip())
            elif not t.startswith("#"):
                cur["bullets"].append(t)
    if cur:
        slides.append(cur)
    if not slides:
        text = md.strip() or "."
        slides = [{"title": fallback, "bullets": [text[:3000]]}]
    return slides


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: gen_pptx.py config.json", file=sys.stderr)
        return 1
    cfg = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
    out = Path(cfg["out"])
    title0 = (cfg.get("title") or out.stem or "Prezentacja").strip()
    md = cfg.get("markdown") or cfg.get("content") or ""

    try:
        from pptx import Presentation
    except ImportError:
        print("Brak biblioteki: pip install python-pptx", file=sys.stderr)
        return 2

    prs = Presentation()
    try:
        layout = prs.slide_layouts[1]
    except IndexError:
        layout = prs.slide_layouts[0]

    for spec in parse_slides(md, title0):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = (spec["title"] or ".")[:255]
        body = slide.shapes.placeholders[1].text_frame
        bullets = spec.get("bullets") or ["."]
        body.text = str(bullets[0])[:4000]
        for b in bullets[1:48]:
            p = body.add_paragraph()
            p.text = str(b)[:4000]
            p.level = 0

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(str(out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
